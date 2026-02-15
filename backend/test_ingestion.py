"""Tests for ingestion endpoints."""
import pytest
from fastapi.testclient import TestClient
from main import app
from io import BytesIO
import os
import time
from test_utils import TEST_EMAIL, TEST_PASSWORD

client = TestClient(app)


def get_auth_token():
    """Get auth token for test user."""
    response = client.post(
        "/auth/login",
        json={"email": TEST_EMAIL, "password": TEST_PASSWORD}
    )
    assert response.status_code == 200, f"Login failed: {response.json()}"
    return response.json()["access_token"]


def test_upload_markdown_file():
    """Test uploading a valid markdown file."""
    # Get auth token
    token = get_auth_token()

    # Create a test markdown file with unique name
    timestamp = int(time.time() * 1000)
    filename = f"test_document_{timestamp}.md"

    md_content = b"""# Test Document

This is a test markdown file for ingestion testing.

## Section 1
Some content here.

## Section 2
More content here.
"""

    # Create file-like object
    files = {
        "file": (filename, BytesIO(md_content), "text/markdown")
    }

    # Upload the file
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )

    # Verify response
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    data = response.json()

    # Verify response data
    assert data["filename"] == filename
    assert data["status"] == "processing"
    assert data["file_size_bytes"] == len(md_content)
    assert "id" in data
    assert "created_at" in data

    print(f"\n[TEST PASSED] Successfully uploaded markdown file: {data['filename']}")
    print(f"  - Document ID: {data['id']}")
    print(f"  - File size: {data['file_size_bytes']} bytes")
    print(f"  - Status: {data['status']}")

    return data["id"]


def test_upload_unsupported_file_type():
    """Test uploading an unsupported file type (PNG)."""
    token = get_auth_token()

    # Create a fake PNG file
    png_content = b"fake png content"
    files = {
        "file": ("screenshot.png", BytesIO(png_content), "image/png")
    }

    # Attempt upload
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )

    # Verify rejection
    assert response.status_code == 400, "Expected 400 for unsupported file type"
    error = response.json()

    # Verify error message includes filename and extension
    assert "screenshot.png" in error["detail"]
    assert ".png" in error["detail"]
    assert "Unsupported file type" in error["detail"]

    print(f"\n[TEST PASSED] Correctly rejected PNG file")
    print(f"  - Error message: {error['detail']}")


def test_upload_oversized_file():
    """Test uploading a file that exceeds size limit."""
    token = get_auth_token()

    # Create a file larger than 10MB
    large_content = b"x" * (11 * 1024 * 1024)  # 11MB
    files = {
        "file": ("large_file.md", BytesIO(large_content), "text/markdown")
    }

    # Attempt upload
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )

    # Verify rejection
    assert response.status_code == 400, "Expected 400 for oversized file"
    error = response.json()

    # Verify error message includes filename and size
    assert "large_file.md" in error["detail"]
    assert "too large" in error["detail"]
    assert "11.0MB" in error["detail"]

    print(f"\n[TEST PASSED] Correctly rejected oversized file")
    print(f"  - Error message: {error['detail']}")


def test_delete_document_cascade():
    """Test that deleting a document cascades to delete chunks automatically."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)

    # Upload document with content to create chunks
    md_content = b"# Test\n\n" + b"Lorem ipsum. " * 200
    files = {"file": (f"test_{timestamp}.md", BytesIO(md_content), "text/markdown")}
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    doc_id = response.json()["id"]

    time.sleep(3)

    # Verify chunks exist
    chunks_response = client.get(
        f"/ingestion/documents/{doc_id}/chunks",
        headers={"Authorization": f"Bearer {token}"}
    )
    assert chunks_response.status_code == 200, f"Failed to get chunks: {chunks_response.json()}"
    chunks = chunks_response.json()
    assert len(chunks) > 0, "Should have chunks before deletion"

    # Delete document
    delete_response = client.delete(
        f"/ingestion/documents/{doc_id}",
        headers={"Authorization": f"Bearer {token}"}
    )
    assert delete_response.status_code == 200, f"Delete failed: {delete_response.json()}"

    # Verify chunks cascade deleted (document no longer exists, so 404)
    chunks_after = client.get(
        f"/ingestion/documents/{doc_id}/chunks",
        headers={"Authorization": f"Bearer {token}"}
    )
    assert chunks_after.status_code == 404, "Chunks should cascade delete"

    print(f"\n[TEST PASSED] Document cascade delete validated")


def test_upload_json_file():
    """Test uploading a valid JSON file."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)
    content = b'{"service":"api","version":"1.0.0","description":"Test JSON document"}'
    files = {"file": (f"test_{timestamp}.json", BytesIO(content), "application/json")}
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    print(f"\n[TEST PASSED] Successfully uploaded JSON file")


def test_upload_csv_file():
    """Test uploading a valid CSV file."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)
    content = b"name,age,city\nAlice,30,NYC\nBob,25,LA\nCharlie,35,Chicago"
    files = {"file": (f"test_{timestamp}.csv", BytesIO(content), "text/csv")}
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    print(f"\n[TEST PASSED] Successfully uploaded CSV file")


def test_upload_xml_file():
    """Test uploading a valid XML file."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)
    content = b'<?xml version="1.0" encoding="UTF-8"?><document><title>Test Document</title><content>Sample XML content for testing.</content></document>'
    files = {"file": (f"test_{timestamp}.xml", BytesIO(content), "application/xml")}
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    print(f"\n[TEST PASSED] Successfully uploaded XML file")


def test_upload_rtf_file():
    """Test uploading a valid RTF file."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)
    content = b'{\\rtf1\\ansi\\deff0 {\\fonttbl {\\f0 Times New Roman;}}\\f0\\fs24 This is a test RTF document with {\\b bold} and {\\i italic} text.}'
    files = {"file": (f"test_{timestamp}.rtf", BytesIO(content), "application/rtf")}
    response = client.post(
        "/ingestion/upload",
        files=files,
        headers={"Authorization": f"Bearer {token}"}
    )
    assert response.status_code == 200, f"Upload failed: {response.json()}"
    print(f"\n[TEST PASSED] Successfully uploaded RTF file")


def test_upload_pptx_file():
    """Test uploading a valid PPTX file."""
    token = get_auth_token()
    timestamp = int(time.time() * 1000)
    try:
        from pptx import Presentation
        from io import BytesIO as BIO
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"
        pptx_bytes = BIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        files = {"file": (f"test_{timestamp}.pptx", pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        response = client.post(
            "/ingestion/upload",
            files=files,
            headers={"Authorization": f"Bearer {token}"}
        )
        assert response.status_code == 200, f"Upload failed: {response.json()}"
        print(f"\n[TEST PASSED] Successfully uploaded PPTX file")
    except ImportError:
        print(f"\n[TEST SKIPPED] PPTX test requires python-pptx library - install with: pip install python-pptx")


if __name__ == "__main__":
    print("=" * 60)
    print("INGESTION UPLOAD TESTS")
    print("=" * 60)

    try:
        # Run tests
        test_upload_markdown_file()
        test_upload_unsupported_file_type()
        test_upload_oversized_file()
        test_delete_document_cascade()
        test_upload_json_file()
        test_upload_csv_file()
        test_upload_xml_file()
        test_upload_rtf_file()
        test_upload_pptx_file()

        print("\n" + "=" * 60)
        print("ALL TESTS PASSED!")
        print("=" * 60)

    except AssertionError as e:
        print(f"\n[TEST FAILED] {e}")
        raise
